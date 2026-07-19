"""DOCX and PPTX rendering for the report exporter.

Both consume the same ``RenderedDocument`` as every other exporter. The optional
libraries (python-docx, python-pptx) are imported lazily inside the render
functions, so importing the exporter package never requires them; the exporter
classes advertise ``available`` from a spec check and raise a clear install
message when missing.

DOCX linearizes each page's elements into native Word content (headings, lists,
tables, images) with real page breaks, headers and footers - editable, not an
image dump. PPTX maps one page to one slide and keeps every element's fractional
position (charts/images as pictures, text as editable text boxes).
"""
from __future__ import annotations

import io
from typing import Any

from fap.reports.layout import RenderedDocument, RenderedElement, RenderedPage

_ROLE_STYLE = {"title": "Title", "subtitle": "Subtitle", "h1": "Heading 1",
               "h2": "Heading 2", "meta": "Normal", "body": "Normal", "caption": "Caption"}


# ================================================================ DOCX
def render_docx(rendered: RenderedDocument, branding: Any = None) -> bytes:
    from docx import Document
    from docx.enum.section import WD_ORIENT, WD_SECTION
    from docx.shared import Emu, Pt

    doc = Document()

    def size_section(section, page: RenderedPage) -> None:
        section.page_width = Emu(int(page.width_pt * 12700))
        section.page_height = Emu(int(page.height_pt * 12700))
        section.orientation = (WD_ORIENT.LANDSCAPE if page.orientation == "landscape"
                               else WD_ORIENT.PORTRAIT)
        t, r, b, l = page.margins_pt
        section.top_margin, section.right_margin = Emu(int(t * 12700)), Emu(int(r * 12700))
        section.bottom_margin, section.left_margin = Emu(int(b * 12700)), Emu(int(l * 12700))

    for i, page in enumerate(rendered.pages):
        if i == 0:
            section = doc.sections[0]
        else:
            section = doc.add_section(WD_SECTION.NEW_PAGE)
        size_section(section, page)
        _docx_hf(section.header, page.header)
        _docx_hf(section.footer, page.footer, number=page.number, confidential=page.confidential)
        for el in _reading_order(page.elements):
            _docx_element(doc, el)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _docx_hf(part, zone, number: str | None = None, confidential: str = "") -> None:
    if zone is None and not number and not confidential:
        return
    p = part.paragraphs[0] if part.paragraphs else part.add_paragraph()
    left = (zone.left if zone else "") or confidential
    center = (zone.center if zone else "")
    right = (zone.right if zone else "") or (number or "")
    # simple three-column header/footer via tabs
    p.text = f"{left}\t{center}\t{right}".rstrip("\t")


def _docx_element(doc, el: RenderedElement) -> None:
    from docx.shared import Inches
    c = el.content
    if el.kind == "spacer":
        doc.add_paragraph("")
        return
    if el.kind == "divider":
        doc.add_paragraph("_" * 40)
        return
    if el.kind in ("image", "chart", "logo"):
        data = c.get("image_bytes")
        if data:
            try:
                doc.add_picture(io.BytesIO(data), width=Inches(6 * el.fw))
            except Exception:
                pass
        if c.get("caption"):
            doc.add_paragraph(c["caption"], style="Caption")
        return
    if el.kind == "kpis":
        rows = c.get("kpis", [])
        table = doc.add_table(rows=1, cols=max(1, len(rows)))
        for j, (k, v) in enumerate(rows):
            cell = table.rows[0].cells[j]
            cell.text = f"{k}\n{v}"
        return
    if el.kind == "table":
        cols, data = c.get("columns", []), c.get("rows", [])
        if cols:
            table = doc.add_table(rows=1, cols=len(cols))
            try:
                table.style = "Light Grid Accent 1"
            except Exception:
                pass
            for j, col in enumerate(cols):
                table.rows[0].cells[j].text = str(col)
            for r in data:
                cells = table.add_row().cells
                for j, val in enumerate(r[:len(cols)]):
                    cells[j].text = str(val)
        return
    if el.kind == "insight":
        doc.add_paragraph(c.get("text", ""), style="Intense Quote")
        return
    # text-like: native headings / bullets / paragraphs
    _docx_text(doc, el)


def _docx_text(doc, el: RenderedElement) -> None:
    style = _ROLE_STYLE.get(el.role, "Normal")
    text = el.content.get("text", "")
    if el.role in ("title", "subtitle") and "\n" not in text:
        try:
            doc.add_paragraph(text, style=style)
        except Exception:
            doc.add_paragraph(text)
        return
    for line in text.splitlines():
        s = line.rstrip()
        if not s:
            continue
        if s.startswith("### "):
            doc.add_paragraph(s[4:], style="Heading 3")
        elif s.startswith("## "):
            doc.add_paragraph(s[3:], style="Heading 2")
        elif s.startswith("# "):
            doc.add_paragraph(s[2:], style="Heading 1")
        elif s.startswith("- "):
            try:
                doc.add_paragraph(s[2:], style="List Bullet")
            except Exception:
                doc.add_paragraph("• " + s[2:])
        else:
            doc.add_paragraph(s)


# ================================================================ PPTX
def render_pptx(rendered: RenderedDocument, branding: Any = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Emu, Pt

    prs = Presentation()
    if rendered.pages:
        p0 = rendered.pages[0]
        prs.slide_width = Emu(int(p0.width_pt * 12700))
        prs.slide_height = Emu(int(p0.height_pt * 12700))
    blank = prs.slide_layouts[6]

    for page in rendered.pages:
        slide = prs.slides.add_slide(blank)
        sw, sh = prs.slide_width, prs.slide_height
        for el in sorted(page.elements, key=lambda e: e.z):
            _pptx_element(slide, el, sw, sh)
        if page.number:
            _pptx_textbox(slide, page.number, 0.9 * sw, 0.94 * sh, 0.08 * sw, 0.05 * sh, size=10)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pptx_element(slide, el: RenderedElement, sw, sh) -> None:
    from pptx.util import Emu
    left, top = Emu(int(el.fx * sw)), Emu(int(el.fy * sh))
    width, height = Emu(int(el.fw * sw)), Emu(int(el.fh * sh))
    c = el.content
    if el.kind in ("image", "chart", "logo"):
        data = c.get("image_bytes")
        if data:
            try:
                slide.shapes.add_picture(io.BytesIO(data), left, top, width, height)
            except Exception:
                pass
        return
    if el.kind in ("spacer", "divider", "cover_overlay"):
        return
    text = _pptx_text_for(el)
    if text:
        _pptx_textbox(slide, text, left, top, width, height,
                      size=_pptx_size(el.role), align=el.align)


def _pptx_text_for(el: RenderedElement) -> str:
    c = el.content
    if el.kind == "kpis":
        return "   ".join(f"{k}: {v}" for k, v in c.get("kpis", []))
    if el.kind == "table":
        cols = c.get("columns", [])
        rows = [" | ".join(map(str, r)) for r in c.get("rows", [])]
        return " | ".join(map(str, cols)) + ("\n" + "\n".join(rows) if rows else "")
    if el.kind == "insight":
        return "‣ " + c.get("text", "")
    return _plain(c.get("text", ""))


def _pptx_textbox(slide, text, left, top, width, height, size=14, align="left") -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Pt
    box = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height)))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in str(text).splitlines() or [""]:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(size)
        para.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        first = False


def _pptx_size(role: str) -> int:
    return {"title": 36, "subtitle": 20, "h1": 22, "h2": 18, "meta": 11,
            "body": 14, "caption": 10}.get(role, 14)


def _reading_order(elements: list[RenderedElement]) -> list[RenderedElement]:
    return sorted(elements, key=lambda e: (round(e.fy, 3), round(e.fx, 3), e.z))


def _plain(text: str) -> str:
    out = []
    for line in (text or "").splitlines():
        s = line.rstrip()
        for m in ("### ", "## ", "# "):
            if s.startswith(m):
                s = s[len(m):]; break
        if s.startswith("- "):
            s = "• " + s[2:]
        out.append(s)
    return "\n".join(out)
